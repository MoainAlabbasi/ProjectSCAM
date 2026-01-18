"""
خدمات الذكاء الاصطناعي - Google Gemini
S-ACM - Smart Academic Content Management System

تم تحديث هذا الملف لاستخدام:
- google-genai SDK الرسمي (الإصدار الجديد)
- موديل gemini-2.5-flash
- Clean Code مع Error Handling محترف

المؤلف: Manus AI
التاريخ: 2026-01-19
"""

from __future__ import annotations

import json
import hashlib
import logging
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from enum import Enum
from functools import wraps
from pathlib import Path
from typing import Optional, List, Dict, Any, Callable, TypeVar, Generic

from django.conf import settings
from django.core.cache import cache

# ========== Logging Configuration ==========
logger = logging.getLogger('ai_features')

# ========== Constants ==========
GEMINI_API_KEY = "AIzaSyC0caScddcPQHxN2fQUSYj02sZ66MG-_80"
GEMINI_MODEL = "gemini-2.5-flash"
MAX_INPUT_LENGTH = 30000
CACHE_TIMEOUT = 3600  # 1 hour
MAX_RETRIES = 3


# ========== Custom Exceptions ==========

class GeminiError(Exception):
    """Base exception for Gemini-related errors."""
    pass


class GeminiConfigurationError(GeminiError):
    """Raised when Gemini is not properly configured."""
    pass


class GeminiAPIError(GeminiError):
    """Raised when Gemini API returns an error."""
    
    def __init__(self, message: str, status_code: Optional[int] = None):
        super().__init__(message)
        self.status_code = status_code


class GeminiRateLimitError(GeminiAPIError):
    """Raised when rate limit is exceeded."""
    pass


class TextExtractionError(GeminiError):
    """Raised when text extraction from file fails."""
    pass


# ========== Enums ==========

class QuestionType(Enum):
    """أنواع الأسئلة المدعومة."""
    MCQ = "mcq"
    TRUE_FALSE = "true_false"
    SHORT_ANSWER = "short_answer"
    MIXED = "mixed"


class ContentType(Enum):
    """أنواع المحتوى المدعومة."""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    TEXT = "text"
    UNKNOWN = "unknown"


# ========== Data Classes ==========

@dataclass
class Question:
    """نموذج السؤال."""
    type: str
    question: str
    answer: str
    options: Optional[List[str]] = None
    explanation: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """تحويل إلى قاموس."""
        result = {
            'type': self.type,
            'question': self.question,
            'answer': self.answer,
        }
        if self.options:
            result['options'] = self.options
        if self.explanation:
            result['explanation'] = self.explanation
        return result


@dataclass
class AIResponse:
    """نموذج استجابة AI."""
    success: bool
    data: Any = None
    error: Optional[str] = None
    cached: bool = False


# ========== Decorators ==========

T = TypeVar('T')


def cache_result(timeout: int = CACHE_TIMEOUT):
    """
    Decorator لتخزين نتائج AI في الكاش.
    
    Args:
        timeout: مدة التخزين بالثواني
        
    Example:
        @cache_result(timeout=3600)
        def generate_summary(self, text: str) -> str:
            ...
    """
    def decorator(func: Callable[..., T]) -> Callable[..., T]:
        @wraps(func)
        def wrapper(self, text: str, *args, **kwargs) -> T:
            # إنشاء مفتاح الكاش
            cache_key = _generate_cache_key(func.__name__, text, args, kwargs)
            
            # محاولة الحصول من الكاش
            cached_result = cache.get(cache_key)
            if cached_result is not None:
                logger.debug(f"Cache hit for {func.__name__}")
                return cached_result
            
            # تنفيذ الدالة
            result = func(self, text, *args, **kwargs)
            
            # تخزين في الكاش
            if result is not None:
                cache.set(cache_key, result, timeout)
                logger.debug(f"Cached result for {func.__name__}")
            
            return result
        return wrapper
    return decorator


def _generate_cache_key(func_name: str, text: str, args: tuple, kwargs: dict) -> str:
    """توليد مفتاح الكاش."""
    content = f"{func_name}:{text}:{str(args)}:{str(sorted(kwargs.items()))}"
    return f"ai:{hashlib.md5(content.encode()).hexdigest()}"


def retry_on_error(max_retries: int = MAX_RETRIES, delay_base: float = 1.0):
    """
    Decorator لإعادة المحاولة عند الفشل.
    
    Args:
        max_retries: عدد المحاولات القصوى
        delay_base: أساس التأخير (exponential backoff)
    """
    def decorator(func: Callable[..., T]) -> Callable[..., T]:
        @wraps(func)
        def wrapper(*args, **kwargs) -> T:
            import time
            last_exception = None
            
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except GeminiRateLimitError as e:
                    last_exception = e
                    delay = delay_base * (2 ** attempt)
                    logger.warning(f"Rate limit hit, retrying in {delay}s (attempt {attempt + 1}/{max_retries})")
                    time.sleep(delay)
                except GeminiAPIError as e:
                    last_exception = e
                    if attempt < max_retries - 1:
                        delay = delay_base * (2 ** attempt)
                        logger.warning(f"API error, retrying in {delay}s: {e}")
                        time.sleep(delay)
                    else:
                        raise
                except Exception as e:
                    logger.error(f"Unexpected error in {func.__name__}: {e}")
                    raise
            
            raise last_exception or GeminiAPIError("Max retries exceeded")
        return wrapper
    return decorator


# ========== Text Extractors ==========

class TextExtractor(ABC):
    """Abstract base class for text extractors."""
    
    @abstractmethod
    def extract(self, file_path: Path) -> str:
        """استخراج النص من الملف."""
        pass
    
    @abstractmethod
    def supports(self, file_path: Path) -> bool:
        """التحقق من دعم نوع الملف."""
        pass


class PDFExtractor(TextExtractor):
    """مستخرج النص من ملفات PDF."""
    
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pdf'
    
    def extract(self, file_path: Path) -> str:
        try:
            import pdfplumber
        except ImportError:
            raise TextExtractionError("pdfplumber not installed. Run: pip install pdfplumber")
        
        try:
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return "\n".join(text_parts)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from PDF: {e}")


class DocxExtractor(TextExtractor):
    """مستخرج النص من ملفات Word."""
    
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.docx'
    
    def extract(self, file_path: Path) -> str:
        try:
            from docx import Document
        except ImportError:
            raise TextExtractionError("python-docx not installed. Run: pip install python-docx")
        
        try:
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from DOCX: {e}")


class PptxExtractor(TextExtractor):
    """مستخرج النص من ملفات PowerPoint."""
    
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pptx'
    
    def extract(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise TextExtractionError("python-pptx not installed. Run: pip install python-pptx")
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from PPTX: {e}")


class PlainTextExtractor(TextExtractor):
    """مستخرج النص من الملفات النصية."""
    
    SUPPORTED_EXTENSIONS = {'.txt', '.md', '.rst', '.csv'}
    
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def extract(self, file_path: Path) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # محاولة مع encoding مختلف
            with open(file_path, 'r', encoding='cp1256') as f:
                return f.read()
        except Exception as e:
            raise TextExtractionError(f"Failed to read text file: {e}")


class TextExtractorFactory:
    """مصنع لإنشاء مستخرجات النص."""
    
    _extractors: List[TextExtractor] = [
        PDFExtractor(),
        DocxExtractor(),
        PptxExtractor(),
        PlainTextExtractor(),
    ]
    
    @classmethod
    def get_extractor(cls, file_path: Path) -> Optional[TextExtractor]:
        """الحصول على المستخرج المناسب للملف."""
        for extractor in cls._extractors:
            if extractor.supports(file_path):
                return extractor
        return None
    
    @classmethod
    def extract_text(cls, file_path: Path) -> str:
        """استخراج النص من الملف."""
        extractor = cls.get_extractor(file_path)
        if extractor is None:
            raise TextExtractionError(f"Unsupported file type: {file_path.suffix}")
        return extractor.extract(file_path)


# ========== Gemini Service ==========

class GeminiService:
    """
    خدمة Google Gemini للذكاء الاصطناعي.
    
    تستخدم google-genai SDK الرسمي مع موديل gemini-2.5-flash.
    
    Features:
        - Clean Code architecture
        - Comprehensive error handling
        - Result caching
        - Retry logic with exponential backoff
        - Text extraction from multiple file formats
    
    Example:
        service = GeminiService()
        
        # توليد تلخيص
        summary = service.generate_summary("نص طويل...")
        
        # توليد أسئلة
        questions = service.generate_questions("نص...", QuestionType.MCQ, 5)
        
        # سؤال المستند
        answer = service.ask_document("نص...", "ما هي الفكرة الرئيسية؟")
    """
    
    def __init__(self, api_key: str = GEMINI_API_KEY, model: str = GEMINI_MODEL):
        """
        تهيئة الخدمة.
        
        Args:
            api_key: مفتاح API لـ Google Gemini
            model: اسم الموديل المستخدم
        """
        self._api_key = api_key
        self._model_name = model
        self._client = None
        
        self._initialize_client()
    
    def _initialize_client(self) -> None:
        """تهيئة عميل Gemini."""
        try:
            from google import genai
            
            self._client = genai.Client(api_key=self._api_key)
            logger.info(f"Gemini client initialized with model: {self._model_name}")
            
        except ImportError:
            raise GeminiConfigurationError(
                "google-genai not installed. Run: pip install google-genai"
            )
        except Exception as e:
            raise GeminiConfigurationError(f"Failed to initialize Gemini client: {e}")
    
    @property
    def is_available(self) -> bool:
        """التحقق من توفر الخدمة."""
        return self._client is not None
    
    def _truncate_text(self, text: str, max_length: int = MAX_INPUT_LENGTH) -> str:
        """قص النص إذا تجاوز الحد الأقصى."""
        if len(text) > max_length:
            logger.warning(f"Text truncated from {len(text)} to {max_length} characters")
            return text[:max_length] + "..."
        return text
    
    @retry_on_error(max_retries=MAX_RETRIES)
    def _generate_content(self, prompt: str, max_tokens: int = 1000) -> str:
        """
        توليد محتوى باستخدام Gemini.
        
        Args:
            prompt: النص المطلوب
            max_tokens: الحد الأقصى للتوكنات
            
        Returns:
            str: النص المولد
            
        Raises:
            GeminiAPIError: عند فشل الـ API
            GeminiRateLimitError: عند تجاوز حد الطلبات
        """
        if not self.is_available:
            raise GeminiConfigurationError("Gemini client not initialized")
        
        try:
            from google import genai
            from google.genai import types
            
            response = self._client.models.generate_content(
                model=self._model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    max_output_tokens=max_tokens,
                    temperature=0.3,
                )
            )
            
            # استخراج النص من الاستجابة
            if response.text:
                return response.text.strip()
            else:
                raise GeminiAPIError("Empty response from Gemini")
                
        except Exception as e:
            error_str = str(e).lower()
            
            if "rate" in error_str or "quota" in error_str:
                raise GeminiRateLimitError(f"Rate limit exceeded: {e}")
            elif "invalid" in error_str and "key" in error_str:
                raise GeminiConfigurationError(f"Invalid API key: {e}")
            else:
                raise GeminiAPIError(f"Gemini API error: {e}")
    
    # ========== Public Methods ==========
    
    def extract_text_from_file(self, file_obj) -> Optional[str]:
        """
        استخراج النص من ملف.
        
        Args:
            file_obj: كائن الملف (LectureFile)
            
        Returns:
            str: النص المستخرج أو None
        """
        if not file_obj.local_file:
            logger.warning(f"File {file_obj.id} has no local file")
            return None
        
        try:
            file_path = Path(file_obj.local_file.path)
            text = TextExtractorFactory.extract_text(file_path)
            logger.info(f"Extracted {len(text)} characters from {file_path.name}")
            return text
        except TextExtractionError as e:
            logger.error(f"Text extraction failed for file {file_obj.id}: {e}")
            return None
    
    @cache_result(timeout=CACHE_TIMEOUT)
    def generate_summary(self, text: str, max_length: int = 500) -> str:
        """
        توليد تلخيص للنص.
        
        Args:
            text: النص المطلوب تلخيصه
            max_length: الحد الأقصى لطول التلخيص
            
        Returns:
            str: التلخيص
            
        Example:
            summary = service.generate_summary("نص طويل جداً...")
        """
        text = self._truncate_text(text)
        
        prompt = f"""أنت مساعد أكاديمي متخصص في تلخيص المحتوى التعليمي باللغة العربية.

قم بتلخيص النص التالي بشكل مختصر ومفيد. ركز على:
- النقاط الرئيسية والمفاهيم الأساسية
- المعلومات الأكثر أهمية
- الحفاظ على الدقة العلمية

النص:
{text}

التلخيص (بحد أقصى {max_length} كلمة):"""

        try:
            return self._generate_content(prompt, max_tokens=max_length * 2)
        except GeminiError as e:
            logger.error(f"Summary generation failed: {e}")
            return self._fallback_summary(text, max_length)
    
    def _fallback_summary(self, text: str, max_length: int) -> str:
        """تلخيص بسيط في حالة فشل الـ AI."""
        sentences = text.replace('\n', ' ').split('.')
        summary = ""
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and len(summary) + len(sentence) < max_length:
                summary += sentence + ". "
            elif len(summary) > 100:
                break
        return summary.strip() or text[:max_length] + "..."
    
    @cache_result(timeout=CACHE_TIMEOUT)
    def generate_questions(
        self, 
        text: str, 
        question_type: QuestionType = QuestionType.MIXED,
        num_questions: int = 5
    ) -> List[Dict[str, Any]]:
        """
        توليد أسئلة من النص.
        
        Args:
            text: النص المصدر
            question_type: نوع الأسئلة
            num_questions: عدد الأسئلة
            
        Returns:
            List[Dict]: قائمة الأسئلة
            
        Example:
            questions = service.generate_questions(
                "نص...", 
                QuestionType.MCQ, 
                5
            )
        """
        text = self._truncate_text(text, 10000)
        
        type_instruction = {
            QuestionType.MCQ: "أسئلة اختيار من متعدد (4 خيارات لكل سؤال)",
            QuestionType.TRUE_FALSE: "أسئلة صح أو خطأ",
            QuestionType.SHORT_ANSWER: "أسئلة إجابة قصيرة",
            QuestionType.MIXED: "مزيج من أنواع الأسئلة المختلفة",
        }.get(question_type, "مزيج من أنواع الأسئلة")
        
        prompt = f"""أنت مدرس متخصص في إنشاء أسئلة اختبارية تعليمية باللغة العربية.

قم بإنشاء {num_questions} سؤال من النص التالي.
نوع الأسئلة المطلوب: {type_instruction}

أرجع الإجابة بصيغة JSON فقط بدون أي نص إضافي:
[
    {{
        "type": "mcq" أو "true_false" أو "short_answer",
        "question": "نص السؤال",
        "options": ["خيار1", "خيار2", "خيار3", "خيار4"],
        "answer": "الإجابة الصحيحة",
        "explanation": "شرح مختصر للإجابة"
    }}
]

ملاحظات:
- للأسئلة من نوع true_false، الخيارات هي ["صح", "خطأ"]
- للأسئلة من نوع short_answer، لا تضع options

النص:
{text}

الأسئلة (JSON فقط):"""

        try:
            result = self._generate_content(prompt, max_tokens=2000)
            return self._parse_questions_json(result)
        except GeminiError as e:
            logger.error(f"Question generation failed: {e}")
            return self._fallback_questions(num_questions)
    
    def _parse_questions_json(self, result: str) -> List[Dict[str, Any]]:
        """تحليل JSON الأسئلة."""
        # تنظيف النتيجة
        result = result.strip()
        
        # إزالة markdown code blocks
        if '```json' in result:
            result = result.split('```json')[1].split('```')[0]
        elif '```' in result:
            result = result.split('```')[1].split('```')[0]
        
        result = result.strip()
        
        try:
            questions = json.loads(result)
            if isinstance(questions, list):
                return questions
            else:
                logger.warning("Questions result is not a list")
                return []
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse questions JSON: {e}")
            logger.debug(f"Raw result: {result[:500]}")
            return []
    
    def _fallback_questions(self, num_questions: int) -> List[Dict[str, Any]]:
        """أسئلة افتراضية في حالة الفشل."""
        return [{
            'type': 'short_answer',
            'question': 'ما هي الفكرة الرئيسية في هذا النص؟',
            'answer': 'راجع النص للإجابة',
            'explanation': 'هذا سؤال تلقائي بسبب عدم توفر خدمة AI'
        }]
    
    def ask_document(self, text: str, question: str) -> str:
        """
        الإجابة على سؤال من سياق المستند.
        
        Args:
            text: نص المستند
            question: السؤال
            
        Returns:
            str: الإجابة
            
        Example:
            answer = service.ask_document(
                "نص المستند...",
                "ما هي الفكرة الرئيسية؟"
            )
        """
        text = self._truncate_text(text)
        
        prompt = f"""أنت مساعد أكاديمي يجيب على الأسئلة بناءً على محتوى المستندات المقدمة.

قواعد الإجابة:
1. أجب بناءً على المحتوى المقدم فقط
2. إذا لم تجد الإجابة في المحتوى، قل ذلك بوضوح
3. استخدم اللغة العربية الفصحى
4. كن واضحاً ومختصراً

المحتوى:
{text}

السؤال: {question}

الإجابة:"""

        try:
            return self._generate_content(prompt, max_tokens=500)
        except GeminiError as e:
            logger.error(f"Document Q&A failed: {e}")
            return "عذراً، حدث خطأ أثناء معالجة سؤالك. يرجى المحاولة مرة أخرى."
    
    def test_connection(self) -> AIResponse:
        """
        اختبار الاتصال بـ Gemini.
        
        Returns:
            AIResponse: نتيجة الاختبار
            
        Example:
            result = service.test_connection()
            if result.success:
                print("Connected!")
        """
        try:
            response = self._generate_content("قل: مرحباً، أنا جاهز!", max_tokens=50)
            return AIResponse(success=True, data=response)
        except GeminiError as e:
            return AIResponse(success=False, error=str(e))


# ========== Celery Tasks (Optional) ==========

try:
    from celery import shared_task
    CELERY_AVAILABLE = True
except ImportError:
    CELERY_AVAILABLE = False
    def shared_task(*args, **kwargs):
        def decorator(func):
            return func
        return decorator


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def generate_summary_async(self, file_id: int) -> Dict[str, Any]:
    """
    مهمة Celery لتوليد التلخيص بشكل غير متزامن.
    
    Args:
        file_id: معرف الملف
        
    Returns:
        Dict: نتيجة التلخيص
    """
    from apps.courses.models import LectureFile
    from apps.ai_features.models import AISummary
    
    try:
        file_obj = LectureFile.objects.get(pk=file_id)
        service = GeminiService()
        
        text = service.extract_text_from_file(file_obj)
        if not text:
            return {'success': False, 'error': 'لا يمكن استخراج النص من الملف'}
        
        summary = service.generate_summary(text)
        
        ai_summary, _ = AISummary.objects.update_or_create(
            file=file_obj,
            defaults={'summary_text': summary, 'is_cached': True}
        )
        
        return {
            'success': True,
            'summary_id': ai_summary.id,
            'summary': summary
        }
        
    except LectureFile.DoesNotExist:
        return {'success': False, 'error': 'الملف غير موجود'}
    except Exception as e:
        logger.error(f"Async summary generation failed: {e}")
        if CELERY_AVAILABLE:
            raise self.retry(exc=e)
        return {'success': False, 'error': str(e)}


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def generate_questions_async(
    self, 
    file_id: int, 
    question_type: str = 'mixed',
    num_questions: int = 5
) -> Dict[str, Any]:
    """
    مهمة Celery لتوليد الأسئلة بشكل غير متزامن.
    
    Args:
        file_id: معرف الملف
        question_type: نوع الأسئلة
        num_questions: عدد الأسئلة
        
    Returns:
        Dict: نتيجة توليد الأسئلة
    """
    from apps.courses.models import LectureFile
    from apps.ai_features.models import AIGeneratedQuestion
    
    try:
        file_obj = LectureFile.objects.get(pk=file_id)
        service = GeminiService()
        
        text = service.extract_text_from_file(file_obj)
        if not text:
            return {'success': False, 'error': 'لا يمكن استخراج النص من الملف'}
        
        q_type = QuestionType(question_type) if question_type in [e.value for e in QuestionType] else QuestionType.MIXED
        questions = service.generate_questions(text, q_type, num_questions)
        
        saved_ids = []
        for q in questions:
            ai_question = AIGeneratedQuestion.objects.create(
                file=file_obj,
                question_type=q.get('type', 'short_answer'),
                question_text=q.get('question', ''),
                options=q.get('options'),
                correct_answer=q.get('answer', ''),
                explanation=q.get('explanation', '')
            )
            saved_ids.append(ai_question.id)
        
        return {
            'success': True,
            'question_ids': saved_ids,
            'count': len(saved_ids)
        }
        
    except LectureFile.DoesNotExist:
        return {'success': False, 'error': 'الملف غير موجود'}
    except Exception as e:
        logger.error(f"Async question generation failed: {e}")
        if CELERY_AVAILABLE:
            raise self.retry(exc=e)
        return {'success': False, 'error': str(e)}
